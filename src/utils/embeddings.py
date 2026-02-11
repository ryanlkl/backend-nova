"""
Docstring for utils.embeddings
"""
# from sentence_transformers import SentenceTransformer
from pypdf import PdfReader
from docx import Document
from io import BytesIO
import csv
import pandas as pd
from pptx import Presentation
import tiktoken
from src.utils.chroma import chroma_client
from datetime import datetime
from src.schema.content import ContentType

# embedding model
# embedding_model = SentenceTransformer("all-MiniLM-L6-v2")

# --- Extract text from content uploads ---
def extract_pdf_text(file_bytes: bytes) -> str:
    reader = PdfReader(BytesIO(file_bytes))
    return "\n".join(page.extract_text() or "" for page in reader.pages)

def extract_docx_text(file_bytes: bytes) -> str:
    doc = Document(BytesIO(file_bytes))
    return "\n".join(p.text for p in doc.paragraphs)

def extract_csv_text(file_bytes: bytes) -> str:
    text = file_bytes.decode("utf-8")
    reader = csv.DictReader(text.splitlines())
    rows = []
    for row in reader:
        rows.append(", ".join(f"{k}: {v}" for k, v in row.items()))
    return "\n".join(rows)

def extract_xlsx_text(file_bytes: bytes) -> str:
    df = pd.read_excel(BytesIO(file_bytes))
    return df.astype(str).to_csv(index=False)

def extract_pptx_text(file_bytes: bytes) -> str:
    prs = Presentation(BytesIO(file_bytes))
    slides = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides.append("\n".join(slide_text))
    return "\n\n".join(slides)

def extract_text(file_bytes: bytes, file_ext: str) -> str:
    if file_ext == "pdf":
        return extract_pdf_text(file_bytes)
    if file_ext == "docx":
        return extract_docx_text(file_bytes)
    if file_ext == "csv":
        return extract_csv_text(file_bytes)
    if file_ext == "xlsx":
        return extract_xlsx_text(file_bytes)
    if file_ext == "pptx":
        return extract_pptx_text(file_bytes)
    raise ValueError(f"Unsupported file type: {file_ext}")

# chunking the text
def chunk_text(
    text: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    model: str = "gpt-4",
):
    encoding = tiktoken.encoding_for_model(model)
    tokens = encoding.encode(text)

    chunks = []
    start = 0
    while start < len(tokens):
        end = start + chunk_size
        chunk_tokens = tokens[start:end]
        chunks.append(encoding.decode(chunk_tokens))
        start += chunk_size - chunk_overlap
    return chunks


def get_collection(name: str):
    return chroma_client.get_or_create_collection(name=name)

def delete_from_vector_db(content_id: str, content_type: ContentType):
    collection = get_collection(content_type.value) # content_type.value

    collection.delete(
        where={"content_id": content_id}
    )



def index_document(
    collection_name: str,
    content_id: str,
    file_ext: str,
    file_bytes: bytes,
    source: str | None = None,
    title: str | None = None,
):
    text = extract_text(file_bytes, file_ext)
    chunks = chunk_text(text, chunk_size=1000, chunk_overlap=200)

    if not chunks:
        return

    documents = []
    metadatas = []
    ids = []

    for i, chunk in enumerate(chunks):
        documents.append(chunk)
        metadatas.append({
            "content_id": content_id,
            "chunk_index": i,
            "source": source,
            "title": title,
            "indexed_at": datetime.utcnow().isoformat(),
        })
        ids.append(f"{content_id}:chunk:{i}")

    collection = get_collection(collection_name)

    collection.upsert(
        documents=documents,
        metadatas=metadatas,
        ids=ids,
    )
    # make test collection tomorrow before merging with existing collection
    # chhaange from test collection to real collection